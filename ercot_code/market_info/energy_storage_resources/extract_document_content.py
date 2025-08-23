#!/usr/bin/env python3
"""
Extract and process content from ERCOT ESR Word and PowerPoint documents.
Creates structured knowledge base from downloaded materials.
"""

import os
import json
import re
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Any
import logging

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

try:
    from docx import Document
    from pptx import Presentation
except ImportError:
    logger.error("Required libraries not installed. Installing now...")
    import subprocess
    subprocess.check_call(["pip", "install", "python-docx", "python-pptx"])
    from docx import Document
    from pptx import Presentation

class DocumentExtractor:
    """Extract content from Word and PowerPoint documents."""
    
    def __init__(self, downloads_dir: str = "downloads"):
        self.downloads_dir = Path(downloads_dir)
        self.extracted_content = {
            "metadata": {
                "extraction_date": datetime.now().isoformat(),
                "source_directory": str(self.downloads_dir)
            },
            "documents": [],
            "ktc_content": {},
            "key_concepts": {},
            "technical_specifications": {},
            "regulatory_timeline": []
        }
    
    def extract_all_documents(self):
        """Process all documents in the downloads directory."""
        logger.info(f"Starting document extraction from {self.downloads_dir}")
        
        # Process each KTC directory
        for ktc_dir in sorted(self.downloads_dir.glob("ktc*")):
            if ktc_dir.is_dir():
                ktc_name = ktc_dir.name
                logger.info(f"Processing {ktc_name}...")
                self.extracted_content["ktc_content"][ktc_name] = {
                    "topic": self._get_ktc_topic(ktc_name),
                    "documents": []
                }
                
                # Process Word documents
                for docx_file in ktc_dir.glob("*.docx"):
                    doc_content = self._extract_docx(docx_file)
                    if doc_content:
                        self.extracted_content["ktc_content"][ktc_name]["documents"].append(doc_content)
                        self.extracted_content["documents"].append(doc_content)
                
                for doc_file in ktc_dir.glob("*.doc"):
                    logger.warning(f"Skipping .doc file (requires conversion): {doc_file.name}")
                
                # Process PowerPoint presentations
                for pptx_file in ktc_dir.glob("*.pptx"):
                    ppt_content = self._extract_pptx(pptx_file)
                    if ppt_content:
                        self.extracted_content["ktc_content"][ktc_name]["documents"].append(ppt_content)
                        self.extracted_content["documents"].append(ppt_content)
        
        # Extract key concepts and technical specifications
        self._extract_key_concepts()
        self._extract_technical_specs()
        self._build_timeline()
        
        logger.info(f"Extraction complete: {len(self.extracted_content['documents'])} documents processed")
    
    def _extract_docx(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from a Word document."""
        try:
            doc = Document(file_path)
            content = {
                "file_name": file_path.name,
                "file_path": str(file_path),
                "file_type": "docx",
                "ktc": file_path.parent.name,
                "content": {
                    "paragraphs": [],
                    "tables": [],
                    "headers": []
                },
                "metadata": {}
            }
            
            # Extract document type from filename
            if "TAC_Approved" in file_path.name:
                content["metadata"]["status"] = "TAC Approved"
                content["metadata"]["approval_date"] = self._extract_date_from_filename(file_path.name)
            elif "Consensus" in file_path.name:
                content["metadata"]["status"] = "BESTF Consensus"
            elif "BESTF" in file_path.name:
                content["metadata"]["status"] = "BESTF Proposal"
            
            # Extract paragraphs
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    # Check for headers
                    if para.style and "Heading" in para.style.name:
                        content["content"]["headers"].append({
                            "level": para.style.name,
                            "text": text
                        })
                    content["content"]["paragraphs"].append(text)
            
            # Extract tables
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    if any(row_data):  # Skip empty rows
                        table_data.append(row_data)
                if table_data:
                    content["content"]["tables"].append(table_data)
            
            logger.info(f"Extracted {file_path.name}: {len(content['content']['paragraphs'])} paragraphs, {len(content['content']['tables'])} tables")
            return content
            
        except Exception as e:
            logger.error(f"Error extracting {file_path}: {e}")
            return None
    
    def _extract_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from a PowerPoint presentation."""
        try:
            prs = Presentation(file_path)
            content = {
                "file_name": file_path.name,
                "file_path": str(file_path),
                "file_type": "pptx",
                "ktc": file_path.parent.name,
                "content": {
                    "slides": [],
                    "titles": [],
                    "bullet_points": []
                },
                "metadata": {
                    "slide_count": len(prs.slides)
                }
            }
            
            # Extract content from each slide
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = {
                    "slide_number": slide_num,
                    "title": "",
                    "content": [],
                    "notes": ""
                }
                
                # Extract title
                if slide.shapes.title:
                    slide_content["title"] = slide.shapes.title.text.strip()
                    content["content"]["titles"].append(slide_content["title"])
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text and text != slide_content["title"]:
                            slide_content["content"].append(text)
                            # Extract bullet points
                            if "\n" in text:
                                bullets = [b.strip() for b in text.split("\n") if b.strip()]
                                content["content"]["bullet_points"].extend(bullets)
                
                # Extract notes
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text:
                        slide_content["notes"] = notes_text.strip()
                
                content["content"]["slides"].append(slide_content)
            
            logger.info(f"Extracted {file_path.name}: {len(content['content']['slides'])} slides")
            return content
            
        except Exception as e:
            logger.error(f"Error extracting {file_path}: {e}")
            return None
    
    def _get_ktc_topic(self, ktc_name: str) -> str:
        """Get the topic for a KTC based on its name."""
        topics = {
            "ktc1": "ESR Registration",
            "ktc2": "PRC and ORDC Reserves",
            "ktc3": "ESR Nodal Pricing and Dispatch/Mitigation",
            "ktc4": "ESR Technical Requirements",
            "ktc5": "Single Model ESRDP",
            "ktc6": "SOC Management and Offer Structure",
            "ktc7": "Base-Point Deviation",
            "ktc8": "Wholesale Storage Load Treatment",
            "ktc9": "Not Available",
            "ktc10": "RUC/COP and Study Assumptions",
            "ktc11": "DC-Coupled Hybrid Resources",
            "ktc12": "AC-Connected Co-Located Resources",
            "ktc13": "Self-Limiting GINRs",
            "ktc14": "New Studies",
            "ktc15": "Proxy Energy Offers and Bids"
        }
        return topics.get(ktc_name, "Unknown Topic")
    
    def _extract_date_from_filename(self, filename: str) -> str:
        """Extract date from filename."""
        # Look for date patterns like 01292020 or 01-29-2020
        patterns = [
            r"(\d{8})",  # 01292020
            r"(\d{2}[-_]\d{2}[-_]\d{4})",  # 01-29-2020
            r"(\d{2}[-_]\d{2}[-_]\d{2})"  # 01-29-20
        ]
        
        for pattern in patterns:
            match = re.search(pattern, filename)
            if match:
                return match.group(1)
        return ""
    
    def _extract_key_concepts(self):
        """Extract key concepts from all documents."""
        concepts = {}
        
        # Define key terms to look for
        key_terms = [
            "Energy Storage Resource", "ESR", "ESRDP", "Battery Energy Storage",
            "State of Charge", "SOC", "Charge", "Discharge", "Cycle",
            "Ancillary Services", "Regulation", "Reserve", "ORDC", "PRC",
            "Registration", "Interconnection", "Technical Requirements",
            "DC-Coupled", "AC-Connected", "Hybrid", "Co-Located",
            "Self-Limiting", "GINR", "Proxy", "Offer", "Bid",
            "Base Point", "Deviation", "Settlement", "Wholesale Storage Load"
        ]
        
        for doc in self.extracted_content["documents"]:
            if doc["file_type"] == "docx":
                text = " ".join(doc["content"]["paragraphs"])
            else:  # pptx
                text = " ".join([" ".join(slide["content"]) for slide in doc["content"]["slides"]])
            
            for term in key_terms:
                if term.lower() in text.lower():
                    if term not in concepts:
                        concepts[term] = {
                            "mentions": 0,
                            "documents": [],
                            "contexts": []
                        }
                    concepts[term]["mentions"] += text.lower().count(term.lower())
                    concepts[term]["documents"].append(doc["file_name"])
                    
                    # Extract context (sentences containing the term)
                    sentences = text.split(".")
                    for sentence in sentences:
                        if term.lower() in sentence.lower() and len(sentence) > 20:
                            concepts[term]["contexts"].append(sentence.strip())
                            break  # Only take first occurrence per document
        
        self.extracted_content["key_concepts"] = concepts
    
    def _extract_technical_specs(self):
        """Extract technical specifications from documents."""
        specs = {
            "registration_requirements": [],
            "technical_parameters": [],
            "operational_constraints": [],
            "market_rules": []
        }
        
        for doc in self.extracted_content["documents"]:
            if "Technical_Requirements" in doc["file_name"] or "Registration" in doc["file_name"]:
                if doc["file_type"] == "docx":
                    # Look for requirements in paragraphs
                    for para in doc["content"]["paragraphs"]:
                        if any(keyword in para.lower() for keyword in ["must", "shall", "require"]):
                            specs["registration_requirements"].append({
                                "requirement": para,
                                "source": doc["file_name"]
                            })
                    
                    # Extract from tables
                    for table in doc["content"]["tables"]:
                        if len(table) > 1:  # Has data rows
                            specs["technical_parameters"].append({
                                "table": table,
                                "source": doc["file_name"]
                            })
        
        self.extracted_content["technical_specifications"] = specs
    
    def _build_timeline(self):
        """Build regulatory timeline from document dates."""
        timeline = []
        
        for doc in self.extracted_content["documents"]:
            if "metadata" in doc and "approval_date" in doc["metadata"]:
                timeline.append({
                    "date": doc["metadata"]["approval_date"],
                    "event": f"{doc['ktc']} - {doc['metadata']['status']}",
                    "document": doc["file_name"]
                })
        
        # Sort by date
        timeline.sort(key=lambda x: x["date"])
        self.extracted_content["regulatory_timeline"] = timeline
    
    def save_results(self, output_file: str = "ercot_esr_extracted_content.json"):
        """Save extracted content to JSON file."""
        with open(output_file, 'w') as f:
            json.dump(self.extracted_content, f, indent=2)
        logger.info(f"Results saved to {output_file}")
    
    def create_knowledge_base(self, output_file: str = "ercot_esr_knowledge_base.json"):
        """Create structured knowledge base from extracted content."""
        knowledge_base = {
            "metadata": {
                "creation_date": datetime.now().isoformat(),
                "total_documents": len(self.extracted_content["documents"]),
                "ktc_topics": len(self.extracted_content["ktc_content"])
            },
            "topics": {},
            "glossary": {},
            "requirements": {},
            "procedures": {},
            "technical_specs": self.extracted_content["technical_specifications"],
            "timeline": self.extracted_content["regulatory_timeline"]
        }
        
        # Organize by topic
        for ktc, ktc_data in self.extracted_content["ktc_content"].items():
            topic = ktc_data["topic"]
            knowledge_base["topics"][topic] = {
                "ktc": ktc,
                "documents": len(ktc_data["documents"]),
                "key_points": [],
                "requirements": [],
                "procedures": []
            }
            
            # Extract key points from documents
            for doc in ktc_data["documents"]:
                if doc["file_type"] == "pptx":
                    # Extract key points from slide titles and bullets
                    for title in doc["content"]["titles"][:5]:  # Top 5 titles
                        knowledge_base["topics"][topic]["key_points"].append(title)
                elif doc["file_type"] == "docx":
                    # Extract from headers
                    for header in doc["content"]["headers"][:5]:
                        knowledge_base["topics"][topic]["key_points"].append(header["text"])
        
        # Build glossary from key concepts
        for concept, data in self.extracted_content["key_concepts"].items():
            if data["contexts"]:
                knowledge_base["glossary"][concept] = {
                    "definition": data["contexts"][0] if data["contexts"] else "",
                    "mentions": data["mentions"],
                    "related_documents": list(set(data["documents"]))[:5]
                }
        
        with open(output_file, 'w') as f:
            json.dump(knowledge_base, f, indent=2)
        logger.info(f"Knowledge base saved to {output_file}")
        
        return knowledge_base


def main():
    """Main execution function."""
    extractor = DocumentExtractor()
    
    # Extract all documents
    extractor.extract_all_documents()
    
    # Save raw extracted content
    extractor.save_results("ercot_esr_extracted_content.json")
    
    # Create knowledge base
    kb = extractor.create_knowledge_base("ercot_esr_knowledge_base.json")
    
    # Print summary
    print("\n" + "="*60)
    print("ERCOT ESR Document Extraction Complete")
    print("="*60)
    print(f"Total documents processed: {len(extractor.extracted_content['documents'])}")
    print(f"KTC topics covered: {len(extractor.extracted_content['ktc_content'])}")
    print(f"Key concepts identified: {len(extractor.extracted_content['key_concepts'])}")
    print(f"Timeline events: {len(extractor.extracted_content['regulatory_timeline'])}")
    print("\nOutput files created:")
    print("  - ercot_esr_extracted_content.json (raw extraction)")
    print("  - ercot_esr_knowledge_base.json (structured knowledge base)")
    
    # Print top concepts
    if extractor.extracted_content['key_concepts']:
        print("\nTop Key Concepts:")
        sorted_concepts = sorted(
            extractor.extracted_content['key_concepts'].items(),
            key=lambda x: x[1]['mentions'],
            reverse=True
        )
        for concept, data in sorted_concepts[:10]:
            print(f"  - {concept}: {data['mentions']} mentions in {len(set(data['documents']))} documents")


if __name__ == "__main__":
    main()
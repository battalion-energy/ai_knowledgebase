#!/usr/bin/env python3
"""
ERCOT ESR/BESS Document Analysis Tool
Extracts and summarizes key information from BESTF Key Topic documents
"""

import os
import re
from pathlib import Path
from docx import Document
from pptx import Presentation
import json
from typing import Dict, List, Tuple
import logging

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class ERCOTDocumentAnalyzer:
    def __init__(self, downloads_dir: str):
        self.downloads_dir = Path(downloads_dir)
        self.extracted_data = {}
        
    def extract_text_from_docx(self, file_path: Path) -> str:
        """Extract text from Word document"""
        try:
            doc = Document(file_path)
            text_content = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text.strip())
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(" | ".join(row_text))
            
            return "\n".join(text_content)
        except Exception as e:
            logger.error(f"Error extracting from {file_path}: {str(e)}")
            return ""
    
    def extract_text_from_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"=== SLIDE {slide_num} ==="]
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                slide_text.append(" | ".join(row_text))
                
                text_content.append("\n".join(slide_text))
            
            return "\n\n".join(text_content)
        except Exception as e:
            logger.error(f"Error extracting from {file_path}: {str(e)}")
            return ""
    
    def extract_key_information(self, text: str, ktc_number: str) -> Dict:
        """Extract key information based on KTC topic"""
        
        # Common patterns to look for
        patterns = {
            'requirements': r'(?i)(requirement|shall|must|required|mandatory)',
            'timelines': r'(?i)(timeline|schedule|date|deadline|implementation|effective)',
            'technical_specs': r'(?i)(MW|kW|ramp rate|SOC|state of charge|frequency|voltage)',
            'testing': r'(?i)(test|testing|validation|verification|compliance)',
            'settlement': r'(?i)(settlement|billing|payment|charge|cost)',
            'telemetry': r'(?i)(telemetry|communication|data|SCADA|EMS)',
        }
        
        extracted = {
            'ktc': ktc_number,
            'key_points': [],
            'requirements': [],
            'timelines': [],
            'technical_specs': [],
            'testing': [],
            'settlement': [],
            'telemetry': []
        }
        
        lines = text.split('\n')
        
        for line in lines:
            line = line.strip()
            if len(line) < 10:  # Skip very short lines
                continue
                
            # Check for different types of information
            for category, pattern in patterns.items():
                if re.search(pattern, line):
                    if category in extracted:
                        extracted[category].append(line)
            
            # Extract bullet points and numbered lists
            if re.match(r'^[\d\w][\.\)]\s', line) or line.startswith('•') or line.startswith('-'):
                extracted['key_points'].append(line)
        
        return extracted
    
    def analyze_ktc_directory(self, ktc_dir: Path) -> Dict:
        """Analyze all documents in a KTC directory"""
        ktc_number = ktc_dir.name.upper()
        logger.info(f"Analyzing {ktc_number}")
        
        ktc_analysis = {
            'ktc': ktc_number,
            'documents': {},
            'summary': {}
        }
        
        for file_path in ktc_dir.iterdir():
            if file_path.is_file():
                logger.info(f"Processing {file_path.name}")
                
                text_content = ""
                if file_path.suffix.lower() == '.docx':
                    text_content = self.extract_text_from_docx(file_path)
                elif file_path.suffix.lower() == '.pptx':
                    text_content = self.extract_text_from_pptx(file_path)
                elif file_path.suffix.lower() == '.doc':
                    # Handle .doc files (legacy format)
                    logger.warning(f"Legacy .doc format detected: {file_path.name}")
                    continue
                
                if text_content:
                    key_info = self.extract_key_information(text_content, ktc_number)
                    ktc_analysis['documents'][file_path.name] = {
                        'text': text_content,
                        'key_info': key_info
                    }
        
        return ktc_analysis
    
    def generate_summary_by_ktc(self) -> Dict:
        """Generate summaries organized by KTC topic"""
        
        ktc_descriptions = {
            'KTC1': 'ESR Registration - Requirements, timelines, and processes',
            'KTC2': 'PRC and ORDC Reserves - Reserve participation rules for ESRs',
            'KTC3': 'Pricing/Dispatch - Nodal pricing, dispatch protocols, mitigation',
            'KTC4': 'Technical Requirements - Interconnection standards, performance specs',
            'KTC5': 'Single Model - Single Model ESRDP implementation',
            'KTC6': 'SOC Management - State of Charge and offer structure requirements',
            'KTC7': 'Base-Point Deviation - BPD rules and penalties',
            'KTC8': 'WSL Treatment - Wholesale Storage Load treatment and settlement',
            'KTC10': 'RUC/COP - RUC and COP submission requirements',
            'KTC11': 'DC-Coupled Hybrid - DC-coupled resource requirements',
            'KTC12': 'AC-Connected Hybrid - AC-connected resource requirements',
            'KTC13': 'Self-limiting GINR - Self-limiting GINR capabilities',
            'KTC15': 'Proxy Bidding - Proxy bidding processes'
        }
        
        summary = {}
        
        for ktc_dir in self.downloads_dir.iterdir():
            if ktc_dir.is_dir() and ktc_dir.name.startswith('ktc'):
                ktc_analysis = self.analyze_ktc_directory(ktc_dir)
                ktc_number = ktc_analysis['ktc']
                
                summary[ktc_number] = {
                    'description': ktc_descriptions.get(ktc_number.upper(), 'Unknown topic'),
                    'analysis': ktc_analysis
                }
        
        return summary
    
    def run_analysis(self) -> Dict:
        """Run the complete analysis"""
        logger.info("Starting ERCOT ESR/BESS document analysis...")
        
        if not self.downloads_dir.exists():
            raise FileNotFoundError(f"Downloads directory not found: {self.downloads_dir}")
        
        return self.generate_summary_by_ktc()

def main():
    downloads_dir = "/pool/ssd8tb/data/iso/ERCOT/market_info/energy_storage_resources/downloads"
    analyzer = ERCOTDocumentAnalyzer(downloads_dir)
    
    try:
        results = analyzer.run_analysis()
        
        # Save results to JSON file
        output_file = "/pool/ssd8tb/data/iso/ERCOT/market_info/energy_storage_resources/ercot_esr_analysis.json"
        with open(output_file, 'w') as f:
            json.dump(results, f, indent=2, default=str)
        
        logger.info(f"Analysis complete. Results saved to {output_file}")
        
        # Print summary
        print("\n" + "="*80)
        print("ERCOT ESR/BESS DOCUMENT ANALYSIS SUMMARY")
        print("="*80)
        
        for ktc, data in results.items():
            print(f"\n{ktc}: {data['description']}")
            print("-" * 60)
            
            analysis = data['analysis']
            print(f"Documents analyzed: {len(analysis['documents'])}")
            
            for doc_name, doc_data in analysis['documents'].items():
                print(f"  - {doc_name}")
                key_info = doc_data['key_info']
                if key_info['key_points']:
                    print(f"    Key points found: {len(key_info['key_points'])}")
                if key_info['requirements']:
                    print(f"    Requirements found: {len(key_info['requirements'])}")
                if key_info['technical_specs']:
                    print(f"    Technical specs found: {len(key_info['technical_specs'])}")
        
        return results
        
    except Exception as e:
        logger.error(f"Analysis failed: {str(e)}")
        raise

if __name__ == "__main__":
    main()